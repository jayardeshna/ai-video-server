import base64
import threading
import shutil

import pkg_resources
import pythoncom
import requests
import google.auth.transport.requests
import cv2
from datetime import datetime

from google.oauth2 import service_account
from pptx import Presentation
from flask import Blueprint, request, jsonify, current_app
from google.cloud import translate_v3beta1 as translate
from gtts import gTTS
from comtypes.client import CreateObject
from concurrent.futures import ThreadPoolExecutor, as_completed
from moviepy.audio.AudioClip import concatenate_audioclips, AudioClip
from moviepy.audio.io.AudioFileClip import AudioFileClip
from moviepy.video.io.VideoFileClip import VideoFileClip
import os
current_time = datetime.now()


file_upload_bp  = Blueprint('file_upload_bp', __name__)

AUDIO_FOLDER = 'uploads/audio'

# JSON data embedded directly in the script
SERVICE_ACCOUNT_INFO = {
  "type": "service_account",
  "project_id": "infra-mechanic-432911-h6",
  "private_key_id": "759f94bde33bbda82b1a20b0952d5cf2892fd080",
  "private_key": "-----BEGIN PRIVATE KEY-----\nMIIEvAIBADANBgkqhkiG9w0BAQEFAASCBKYwggSiAgEAAoIBAQDI1AOiCijV5IYm\npGS0KDHa9g/GHqYQDztGf5Zvr/0lh26At4IrjdZnoj8Wj/pOvcOd5aBFaxPIf3w9\nbYeL5qDDPVL/Ww/oB4NvM7UPtQImtGMoJxEicknZiqeL6/AiGoUTlSZMXp28k6lk\nBOW3AYwauag0mtSMY3kzzq9CHs66VYGZvUl8lfAEqv/nHWlBj8/PvtKDkOIDw9JO\n2Ak2z2MbZvf9CU87RA2bzq4t4fDsXZtB7P5iZ8qLGexJ4XPANUXl4zIm6y3l3RnP\n8kwqBnKhQO0XORuiN53EJiMhWD37ry81TrQNvtE2V9mrllJB0xmwNLSjzwjdykrL\nFARCUaYDAgMBAAECggEAEMm8ZgLQaBT66+cJLhBAHqw8JFUyj5weWeXXhw2grU71\ne1gbzxVDy5UCDhpKE+QEIsJr0/HADR4vHhl2X8kPPlL1ibgsi7p5D8Y0SBRbEi45\nLcWaRpVkPAW9vNIj2E38OudYdMDiEb4MRnqsU1CTTDah64fX5cZ4/s4z02Ss5mGS\nk08xOsEBRW0VCahl4VKg8dtpg+eSd5wjLr7mkTMngGUnwzJrf7sZ9qVmIv1PlrEe\n+KlSADzsE7Zn8i4JNnDk+EF79VCiq72OMxCnGIj0CGG1Uk3lHMDxg2bmYJvQ9mOP\nOAwUWpFtEGy8cWG6/BpmN0tzPDrTIiP/EXryBN8NMQKBgQDkwKll4rdeJLafIqGH\nCFupMdTMvBNIQgXPziVE5avCCv5OpTGsnHbQvNLqw+4qaM1VRJKgdSrvLVqoowYi\nq3VvpQluodh1OsTvYJKW4ZCSELKm3sFTlutHOIU48YxT6SLqHU/+Un5AUld+Fn9n\nnEQwT700o/KxzjcwWEmShFt+swKBgQDgv9r471j5E7x4MMGHuWDnZPWSBNcvNUvE\n66ezd34b1wT1ri/3MPBwNrLUtV+2uD1evgvPMWTvVg+nD9UPZ9+ptuNVQDcBo91E\njOm2TWNSrrUehkKV2edqjWKLBkUBxwLKbWxiqlLUm01BStSXnl4W1Y9IgFw3zSUA\nieurBBLjcQKBgDJkMqzsHaLSiZxSv6yEbdS+0nbrc3/2c1tMuS8NPSH40/61K1Uv\n/oiLIBdxY/TZs2n2Oiq4xOV9YzxE9WQLy0n13DcP1iQ09w8YghJUzEkw/nn+Fi54\nz4mX7NkXdOdDM1rJbEbFG2TF/toa0KQda3QNDMwc1ajPmZ6/xC1PlOw7AoGAAaAh\nm09P221THaNsouiqVMIcLJQwYj0GcUQrEFc66ESj2osJ+cW213jtzjAJoy6t8vN4\n7mSzEDAVKOYqalQcAAQ8nK88NpzSyt4iCpmAsGZnFPNvo9nRnKmtHshK03p3ALPN\n5uDIfJUq99/srbQNhijBQor7H8QyElDtDW2YHWECgYBetrVfTMB70jVhCbtZ32Pl\nWWhxra1zF7sGM9fuL85qQ2RbK3jBLPwnrZxcz3Y2F4pU/ZxtrDdLrb4noJSjCpP0\n+F4Y9RvwgXFcA21vI65NPsDEbjPAh8Dp6Zswdt1tRxkdMSMK5hI/fdFiaeLT8HSG\nnI/yOrEA279JAA8sE8NF6Q==\n-----END PRIVATE KEY-----\n",
  "client_email": "translation-service-account@infra-mechanic-432911-h6.iam.gserviceaccount.com",
  "client_id": "108658361032202589941",
  "auth_uri": "https://accounts.google.com/o/oauth2/auth",
  "token_uri": "https://oauth2.googleapis.com/token",
  "auth_provider_x509_cert_url": "https://www.googleapis.com/oauth2/v1/certs",
  "client_x509_cert_url": "https://www.googleapis.com/robot/v1/metadata/x509/translation-service-account%40infra-mechanic-432911-h6.iam.gserviceaccount.com",
  "universe_domain": "googleapis.com"
}


@file_upload_bp.route('/api/v1/upload', methods=['POST'])
def upload_ppt():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400

    file = request.files['file']

    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    upload_folder = current_app.config['PPT_FOLDER']

    # Check if the directory exists, and create it if it doesn't
    if not os.path.exists(upload_folder):
        os.makedirs(upload_folder)
    print("Starting here:", datetime.now())
    if file and file.filename.endswith('.pptx'):
        filepath = os.path.join(current_app.config['PPT_FOLDER'], file.filename)
        file.save(filepath)
        translated_filepath = os.path.join(current_app.config['PPT_FOLDER'], 'Translated_' + file.filename)
        target_language = request.args.get('language', 'Hindi')
        language_code = request.args.get('language_code', '')
        print("Transalting started in google")
        translate_ppt1(filepath, translated_filepath, language_code)
        slides_data = extract_text_from_ppt(translated_filepath)
        image_folder = 'uploads/images/'
        thread = threading.Thread(target=create_audio_from_text, args=(slides_data, language_code))
        thread.start()
        save_presentation_as_images(translated_filepath, image_folder)
        translated_slides_data = []
        for slide in slides_data:
            combined_text = " ".join(slide['texts'])
            translated_slides_data.append({
                'slide_number': slide['slide_number'],
                'texts': combined_text,
            })
        return jsonify({
            'message': 'File uploaded and text extracted successfully',
            'filename': 'Translated_' + file.filename,
            'slides': translated_slides_data
        }), 200
    return jsonify({'error': 'Invalid file type'}), 400


@file_upload_bp.route('/api/v1/generate', methods=['POST'])
def generate_video():
    print("hello")
    data = request.get_json()
    slides = data.get('slides', [])
    file_name = data.get('filename')
    language_code = request.args.get('language_code', '')
    transition_delay_str = request.args.get('transition_delay')
    try:
        tansition_delay = int(transition_delay_str)
    except ValueError:
        # Handle the case where conversion fails
        tansition_delay = None  # or any default value you prefer
    ppt_path = os.path.join('uploads/ppt/', file_name)

    audio_folder = 'uploads/audio'
    video_folder = 'uploads/videos'
    image_folder = 'uploads/images/'
    os.makedirs(image_folder, exist_ok=True)
    os.makedirs(audio_folder, exist_ok=True)
    os.makedirs(video_folder, exist_ok=True)


    video_file = os.path.join(video_folder, "output_video.mp4")

    # Set the frame size (width and height) and fps
    first_image = cv2.imread(slides_images[0])
    height, width, layers = first_image.shape
    frame_size = (width, height)
    fps = 1  # Adjust this if needed

    fourcc = cv2.VideoWriter_fourcc(*'mp4v')
    out = cv2.VideoWriter(video_file, fourcc, fps, frame_size)

    audio_clips = []
    for slide in slides:
        slide_number = slide.get('slide_number')
        slide_image_path = slides_images[slide_number - 1]

        audio_file = f"slide_{slide_number}.mp3"
        audio_path = os.path.join(audio_folder, audio_file)
        audio_clip = AudioFileClip(audio_path)
        audio_clips.append(audio_clip)

        # Write the image to the video
        frame_count = int((audio_clip.duration + 1.5) * fps)
        print(frame_count, "frame_count")
        image = cv2.imread(slide_image_path)  # Read the image into a NumPy array
        for _ in range(frame_count):
            out.write(image)
    out.release()
    silence = create_silence(1.5)
    results = []
    for clip in audio_clips:
        result = process_clip(clip, silence)
        results.append(result)
    # Flatten the list of results
    audio_clips_with_delay = [item for sublist in results for item in sublist]
    # Create final audio
    final_audio = concatenate_audioclips(audio_clips_with_delay)

    # Load the video and add audio
    video = VideoFileClip(video_file)
    final_video = video.set_audio(final_audio)

    final_video_path = os.path.join('uploads/', file_name + '.mp4')
    final_video.write_videofile(final_video_path,
                                codec='libx264',
                                audio_codec='aac',
                                bitrate='1000k',
                                threads=4,  # Increase number of threads
                                preset='veryfast',  # Change preset to a faster option
                                ffmpeg_params=['-crf', '23'])
    shutil.rmtree(audio_folder, ignore_errors=True)
    shutil.rmtree(image_folder, ignore_errors=True)
    # shutil.rmtree(current_app.config['PPT_FOLDER'], ignore_errors=True)
    shutil.rmtree(video_folder, ignore_errors=True)

    return jsonify(
        {'message': 'Presentation updated successfully', 'final_video': os.path.abspath(final_video_path)}), 200


def create_audio_from_text(slides_data, language_code):
    print(datetime.now(), "starting of creating audio")
    # Define the worker function
    def process_slide(slide):
        text = slide['texts']
        if isinstance(text, list):
            text = " ".join(text)
        index = slide['slide_number']
        audio_folder = 'uploads/audio'
        os.makedirs(audio_folder, exist_ok=True)
        audio_file = f"slide_{index}.mp3"
        audio_path = os.path.join(audio_folder, audio_file)
        text_to_speech(text, language_code, output_file=audio_path)
        return audio_path

    # Use ThreadPoolExecutor to process each slide in a separate thread
    with ThreadPoolExecutor() as executor:
        futures = {executor.submit(process_slide, slide): slide['slide_number'] for slide in slides_data}
        for future in as_completed(futures):
            slide_number = futures[future]
            try:
                audio_path = future.result()
                print(f"Slide {slide_number} audio saved at {audio_path}")
            except Exception as e:
                print(f"Error processing slide {slide_number}: {e}")


def get_access_token(credentials):
    auth_request = google.auth.transport.requests.Request()
    credentials.refresh(auth_request)
    return credentials.token

def get_service_account_credentials():
    return service_account.Credentials.from_service_account_info(SERVICE_ACCOUNT_INFO)


def translate_ppt1(input_ppt_path, output_file_path, target_language_code):
    credentials = get_service_account_credentials()
    client = translate.TranslationServiceClient(credentials=credentials)
    project_id = "infra-mechanic-432911-h6"
    location = "us-central1"
    parent = f"projects/{project_id}/locations/{location}"
    with open(input_ppt_path, "rb") as document:
        document_content = document.read()
    document_input_config = {
        "content": document_content,
        "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }

    response = client.translate_document(
        request={
            "parent": parent,
            "target_language_code": target_language_code,
            "document_input_config": document_input_config,
        }
    )
    translated_content = response.document_translation.byte_stream_outputs[0]

    with open(output_file_path, "wb") as output_file:
        output_file.write(translated_content)
    print(f"Document translated successfully! Saved as '{output_file_path}'.")
    print(f"Detected Language Code: {response.document_translation.detected_language_code}")

# TODO Here is the code for translating ppt using rest API.

# def translate_ppt(input_ppt_path, output_file_path, target_language_code):
#     scopes = ["https://www.googleapis.com/auth/cloud-translation"]
#     credentials = service_account.Credentials.from_service_account_info(SERVICE_ACCOUNT_INFO, scopes=scopes)
#     access_token = get_access_token(credentials)
#
#     url = f"https://translate.googleapis.com/v3beta1/projects/infra-mechanic-432911-h6/locations/us-central1:translateDocument"
#
#     headers = {
#         "Authorization": f"Bearer {access_token}",
#         "Content-Type": "application/json"
#     }
#
#     with open(input_ppt_path, "rb") as document:
#         document_content = document.read()
#     encoded_content = base64.b64encode(document_content).decode('utf-8')
#
#     body = {
#         "targetLanguageCode": target_language_code,
#         "documentInputConfig": {
#             "content": encoded_content,  # Decode to string to include in JSON
#             "mimeType": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
#         }
#     }
#
#     response = requests.post(url, headers=headers, json=body)
#
#     if response.status_code == 200:
#         result = response.json()
#         translated_content = base64.b64decode(result['documentTranslation']['byteStreamOutputs'][0])
#         with open(output_file_path, "wb") as output_file:
#             output_file.write(translated_content)
#
#         print(f"Document translated successfully! Saved as '{output_file_path}'.")
#         print(f"Detected Language Code: {result['documentTranslation']['detectedLanguageCode']}")
#     else:
#         print(f"Error: {response.status_code} - {response.text}")


def create_silence(duration):
    return AudioClip(lambda t: 0, duration=duration)


def process_clip(clip, silence):
    return [clip, silence]


slides_images =[]

def save_presentation_as_images(ppt_path, output_folder):
    pythoncom.CoInitialize()
    powerpoint = CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1

    ppt_abs_path = os.path.abspath(ppt_path)
    output_folder_abs = os.path.abspath(output_folder)
    # Open the presentation
    presentation = powerpoint.Presentations.Open(ppt_abs_path)

    if not os.path.exists(output_folder_abs):
        os.makedirs(output_folder_abs)

    for i, slide in enumerate(presentation.Slides):
        slide_image_path = os.path.join(output_folder_abs, f"slide_{i + 1}.jpg")
        slides_images.append(slide_image_path)
        slide.Export(slide_image_path, "JPG", 1920, 1080)  # You can adjust resolution here
        print(f"Saved {slide_image_path}")

    # Close the presentation and quit PowerPoint
    presentation.Close()
    powerpoint.Quit()
    pythoncom.CoUninitialize()
    return slides_images


def extract_text_from_ppt(filepath):
    prs = Presentation(filepath)
    slides_data = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_content = {
            'slide_number': slide_number,
            'texts': []
        }

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                # cleaned_text = clean_text(shape.text)
                slide_content['texts'].append(shape.text)

        slides_data.append(slide_content)

    return slides_data


def text_to_speech(text, lang, output_file):
    if not text or not text.strip():  # Check if text is None, empty, or only whitespace
        text = "This image is displayed for a few seconds. Please pause the video if you need more time to read or understand the content."
    tts = gTTS(text=text, lang=lang, slow=False)
    tts.save(output_file)
    return output_file